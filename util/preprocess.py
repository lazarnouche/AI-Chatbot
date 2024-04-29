import pickle
import fitz  
from PIL import Image
from io import BytesIO
import os
import numpy as np
import json
from langchain_community.document_loaders import PyPDFLoader

def split_into_chuncks(docs,chunk_size=3000):
  splitted_docs = []
  for doc in docs :
    chunks = split_text(doc, chunk_size)
    splitted_docs.extend(chunks)
  return splitted_docs

def split_text(long_text, chunk_size):
  chunks = []
  for i in range(0, len(long_text), chunk_size):
      chunk = long_text[i:i + chunk_size]
      chunks.append(chunk)
  return chunks

def read_json(path):
    with open(path) as user_file:
        file_contents = user_file.read()
    return json.loads(file_contents)

def save_file(file_path, text):
  with open(file_path, 'wb') as file:
    # A new file will be created
    pickle.dump(text, file)

def load_file(file_path):
  with open(file_path, 'rb') as file:
      # dump information to that file
      file = pickle.load(file)
      return file

def pdf_page_to_pillow_image(pdf_path, page_number=0):
    
  doc = fitz.open(pdf_path)
  
  if page_number < 0 or page_number >= doc.page_count:
      raise ValueError("Invalid page number")
  page = doc[page_number]
  image_bytes = page.get_pixmap().tobytes()
  
  doc.close()
  # Create a Pillow Image from the image bytes
  image = Image.open(BytesIO(image_bytes))
  return image

def load_pdf(file_path):
    loader = PyPDFLoader(file_path)
    pages = loader.load_and_split()
    return pages

def clean_document_path(un_processed_docs, path_to_change = ""):
    from langchain.docstore.document import Document
    docs = []
    for pg in un_processed_docs:
        dic_page = dict(pg)
        dic_page['metadata']['source'] = path_to_change
        doc =  Document(page_content = dic_page['page_content'],
                        metadata=dic_page['metadata'] )
        docs.append(doc)
    return docs

def read_images_in_pdf(pages_dic,n_images):
  imgs = []
  for i,page_dic in enumerate(pages_dic):
    filepath_pdf = page_dic['source']
    page = page_dic['page']
    if not os.path.isfile(filepath_pdf):
      continue
    img = pdf_page_to_pillow_image(filepath_pdf, page)
    if i+1>n_images:
      break
    imgs.append(img)
  return imgs
  
def read_ppt_and_convert_to_document(file_path):
    from pptx import Presentation
    from langchain.docstore.document import Document
    presentation = Presentation(file_path)
    pages = []
    for i, slide in enumerate(presentation.slides):
        # Access text content in each shape on the slide
        page = ''
        position_top, position_left, list_page=[],[],[]
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                position_top.append(shape.top)
                position_left.append(shape.left)
                list_page.append(shape.text)
        indices_top = list(np.argsort(position_top))
        index_min_left = np.argmin(position_left)
        element = indices_top.pop(indices_top.index(index_min_left))
        indices_top.insert(1, element)
        for j in indices_top:
            page+= ' '+ list_page[j]
        doc =  Document(page_content=page, metadata={"source": file_path,"page": i})
        pages.append(doc)
    return pages